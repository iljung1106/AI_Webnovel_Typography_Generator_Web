from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DIR = ROOT / "docs" / "reports"
PPTX_OUT = OUT_DIR / "웹소설_타이포그래피_중간발표_14장.pptx"
TEMPIMG = OUT_DIR / "tempimg"

BLUE = RGBColor(37, 99, 235)
NAVY = RGBColor(24, 39, 74)
INK = RGBColor(17, 24, 39)
MUTED = RGBColor(100, 116, 139)
LINE = RGBColor(213, 223, 235)
SOFT = RGBColor(244, 248, 252)
GREEN = RGBColor(20, 150, 120)
PURPLE = RGBColor(111, 91, 232)
RED = RGBColor(211, 75, 75)
AMBER = RGBColor(196, 128, 28)


SCREENSHOTS = {
    "genre": TEMPIMG / "스크린샷 2026-05-25 232011.png",
    "cover": TEMPIMG / "스크린샷 2026-05-25 232403.png",
    "title": TEMPIMG / "스크린샷 2026-05-25 232415.png",
    "layout": TEMPIMG / "스크린샷 2026-05-25 232426.png",
    "style": TEMPIMG / "스크린샷 2026-05-25 232437.png",
    "generation": TEMPIMG / "스크린샷 2026-05-25 232933.png",
    "effect": TEMPIMG / "스크린샷 2026-05-25 234002.png",
    "export": TEMPIMG / "스크린샷 2026-05-25 234024.png",
}


def emu(v: float):
    return Inches(v)


def set_font(run, size=18, bold=False, color=INK):
    run.font.name = "Malgun Gothic"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def text_box(slide, x, y, w, h, text="", size=18, bold=False, color=INK, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    for r in p.runs:
        set_font(r, size=size, bold=bold, color=color)
    return shape


def title(slide, no: str, title_text: str, sub: str | None = None):
    text_box(slide, 0.62, 0.42, 2.4, 0.24, no, size=8.5, bold=True, color=BLUE)
    text_box(slide, 0.62, 0.75, 8.8, 0.54, title_text, size=25, bold=True, color=INK)
    if sub:
        text_box(slide, 0.64, 1.35, 10.4, 0.3, sub, size=10.5, color=MUTED)


def footer(slide, n: int):
    text_box(slide, 0.62, 7.08, 4.0, 0.22, "서비스데이터사이언스 중간발표", size=7.4, color=MUTED)
    text_box(slide, 12.25, 7.08, 0.5, 0.22, f"{n:02d}", size=7.4, bold=True, color=BLUE, align=PP_ALIGN.RIGHT)


def rect(slide, x, y, w, h, fill=SOFT, line=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    s = slide.shapes.add_shape(shape_type, emu(x), emu(y), emu(w), emu(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(0.7)
    return s


def add_image(slide, path: Path, x, y, w, h, border=True):
    with Image.open(path) as im:
        iw, ih = im.size
    scale = min(w / iw, h / ih)
    pw, ph = iw * scale, ih * scale
    px, py = x + (w - pw) / 2, y + (h - ph) / 2
    if border:
        rect(slide, x, y, w, h, fill=RGBColor(255, 255, 255), line=LINE)
    slide.shapes.add_picture(str(path), emu(px), emu(py), width=emu(pw), height=emu(ph))


def bullet_lines(slide, x, y, lines, size=13.5, gap=0.43, color=INK, dot_color=BLUE):
    for i, line in enumerate(lines):
        cy = y + i * gap
        d = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(x), emu(cy + 0.1), emu(0.07), emu(0.07))
        d.fill.solid()
        d.fill.fore_color.rgb = dot_color
        d.line.color.rgb = dot_color
        text_box(slide, x + 0.18, cy, 5.6, 0.3, line, size=size, color=color)


def pill(slide, x, y, w, h, head, body="", accent=BLUE, fill=RGBColor(255, 255, 255)):
    s = rect(slide, x, y, w, h, fill=fill, line=RGBColor(198, 213, 240))
    tf = s.text_frame
    tf.clear()
    tf.margin_left = emu(0.1)
    tf.margin_right = emu(0.1)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = head
    set_font(r, size=12, bold=True, color=accent)
    if body:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = body
        set_font(r2, size=7.6, color=MUTED)
    return s


def arrow(slide, x1, y1, x2, y2, color=BLUE):
    s = slide.shapes.add_connector(1, emu(x1), emu(y1), emu(x2), emu(y2))
    s.line.color.rgb = color
    s.line.width = Pt(1.5)
    s.line.end_arrowhead = True


def cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    text_box(s, 0.62, 0.55, 2.5, 0.26, "AI 기반 웹소설 타이포그래피", size=9, bold=True, color=BLUE)
    text_box(s, 0.62, 1.08, 6.3, 1.35, "AI 기반 웹소설\n타이포그래피 제작 웹앱", size=34, bold=True)
    text_box(s, 0.68, 2.76, 5.8, 0.62, "생성형 AI와 웹 프레임워크를 결합해\n누구나 손쉽게 표지 제목 타이포그래피를 만드는 서비스", size=13, color=MUTED)
    rect(s, 0.62, 3.7, 1.45, 0.32, fill=RGBColor(235, 241, 255), line=RGBColor(197, 211, 255))
    text_box(s, 0.72, 3.76, 1.25, 0.18, "중간 발표", size=8.5, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    add_image(s, SCREENSHOTS["effect"], 7.05, 0.82, 5.6, 4.35)
    text_box(s, 7.15, 5.45, 5.35, 0.25, "실제 구현 화면: 효과 편집 단계", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    rect(s, 0.62, 6.32, 12.1, 0.02, fill=BLUE, line=BLUE, radius=False)
    text_box(s, 0.62, 6.5, 3.4, 0.24, "발표자: 김일중", size=8.6, color=MUTED)
    text_box(s, 9.2, 6.5, 3.5, 0.24, "2026.05.26", size=8.6, color=MUTED, align=PP_ALIGN.RIGHT)


def problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "02 Problem", "웹소설 시장의 성장, 그러나 타이포 제작은 여전히 무겁다")
    text_box(s, 0.85, 1.75, 11.4, 0.48, "제목 타이포그래피는 작품의 첫인상을 결정하지만, 개인 작가와 소규모 제작자에게는 비용과 시간이 큰 장벽입니다.", size=18, bold=True, align=PP_ALIGN.CENTER)
    items = [("비용", "타이포 외주 비용 부담\n수십만 원 단위까지 상승"), ("시간", "시안·수정에 며칠~몇 주\n빠른 테스트에 부적합"), ("도구", "일반 편집 툴은 웹소설\n장르 문법에 최적화되지 않음")]
    for i, (h, b) in enumerate(items):
        pill(s, 1.0 + i * 4.05, 3.0, 3.15, 1.15, h, b, accent=[RED, AMBER, BLUE][i])
    text_box(s, 1.0, 5.3, 11.3, 0.42, "직접 만들기 어렵고, 외주는 비싸고, 기존 툴은 웹소설 제목 타이포 제작에 초점이 맞지 않습니다.", size=17, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 2)


def product_goal(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "03 Product Goal", "외주를 대체하기보다, 외주 전 단계의 시안을 빠르게")
    text_box(s, 0.95, 1.82, 11.2, 0.9, "정식 외주 전 단계에서\n작품 제목과 장르를 바탕으로 고품질 타이포그래피 시안을 빠르게 생성·수정·내보내는 것", size=24, bold=True, align=PP_ALIGN.CENTER)
    pill(s, 1.25, 3.42, 2.6, 0.8, "외주 비용 절감", "초기 시안 비용 부담 완화", accent=GREEN)
    pill(s, 4.15, 3.42, 2.6, 0.8, "3분 내 시안", "빠른 방향성 검토", accent=BLUE)
    pill(s, 7.05, 3.42, 2.6, 0.8, "장르 문법 반영", "웹소설식 제목 인상", accent=PURPLE)
    pill(s, 9.95, 3.42, 2.0, 0.8, "편집 가능", "효과·배치 조정", accent=NAVY)
    text_box(s, 1.0, 5.45, 11.2, 0.4, "미션: 매력적인 제목 디자인을 더 많은 작가가 시도할 수 있게 만드는 것", size=17, bold=True, color=INK, align=PP_ALIGN.CENTER)
    footer(s, 3)


def competitive(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "04 Competitive", "범용 이미지 생성기는 ‘표지용 타이포 파일’까지 해결하지 못한다")
    text_box(s, 0.75, 1.75, 12.0, 0.35, "ChatGPT Images 2.0, Google Nano Banana 2 같은 최신 이미지 생성기도 강력하지만, 웹소설 제목 타이포 제작에는 별도 요구사항이 있습니다.", size=13.8, bold=True, align=PP_ALIGN.CENTER)
    left = ["한글 글자 구조와 가독성 유지", "흑백 실루엣 기반 중간 산출물", "투명 배경 PNG", "효과·레이어 분리 export", "표지 위 실시간 배치 수정"]
    right = ["이미지 전체 생성에 강함", "텍스트 렌더링은 개선 중", "투명·레이어 산출물이 기본 목표는 아님", "웹소설 장르별 타이포 문법은 별도 설계 필요"]
    text_box(s, 1.05, 2.75, 4.8, 0.3, "우리 서비스가 필요한 요구사항", size=15, bold=True, color=BLUE)
    bullet_lines(s, 1.05, 3.18, left, size=12.3)
    text_box(s, 7.05, 2.75, 4.7, 0.3, "범용 생성기의 한계", size=15, bold=True, color=RED)
    bullet_lines(s, 7.05, 3.18, right, size=12.3, dot_color=RED)
    footer(s, 4)


def framing(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "05 AI/ML Framing", "한두 줄의 자연어를 전문적인 타이포 제작 단계로 변환")
    text_box(s, 0.9, 1.82, 11.5, 0.38, "Input: 작품 제목 + 장르 + 분위기/요소 → Output: 타이포그래피 후보와 표지 합성 결과", size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    stages = [("Title", "작품 제목"), ("Genre", "장르"), ("Mood", "분위기·요소"), ("LLM", "디자인 프롬프트"), ("Diffusion", "흑백 타이포"), ("Edit", "효과·export")]
    for i, (h, b) in enumerate(stages):
        x = 0.65 + i * 2.1
        pill(s, x, 3.0, 1.65, 0.86, h, b, accent=BLUE if i < 3 else GREEN if i == 3 else PURPLE)
        if i < len(stages) - 1:
            arrow(s, x + 1.67, 3.43, x + 2.02, 3.43)
    text_box(s, 1.1, 5.35, 11.0, 0.55, "거시적으로는 자연어 → 타이포그래피 생성 구조이고,\n세부적으로는 전문가 작업 프로세스를 단계별로 나눈 구조입니다.", size=18, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 5)


def pipeline(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "06 AI Pipeline", "생성형 AI에 전부 맡기지 않는 Hybrid Approach")
    text_box(s, 0.8, 1.72, 11.6, 0.35, "비용과 가독성 문제를 줄이기 위해, 생성과 편집을 역할별로 분리합니다.", size=15, bold=True, align=PP_ALIGN.CENTER)
    items = [("LLM", "장르·요소를\n프롬프트로 정리"), ("Diffusion", "흑백 실루엣\n타이포 생성"), ("Post-processing", "투명 PNG\n마스크 변환"), ("Frontend Render", "색·재질·빛\n실시간 효과"), ("Export", "PNG / ZIP\n결과 저장")]
    for i, (h, b) in enumerate(items):
        x = 0.95 + i * 2.35
        pill(s, x, 3.02, 1.85, 0.92, h, b, accent=[BLUE, PURPLE, GREEN, NAVY, AMBER][i])
        if i < len(items) - 1:
            arrow(s, x + 1.88, 3.48, x + 2.28, 3.48)
    text_box(s, 1.0, 5.38, 11.2, 0.6, "핵심은 흑백 뼈대만 AI로 만들고, 컬러·질감·빛 효과는 웹 렌더링으로 처리해\n서버 부담과 재생성 비용을 줄이는 것입니다.", size=17, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 6)


def user_flow(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "07 User Flow", "복잡한 프롬프트 대신 단 몇 번의 클릭으로")
    steps = ["로그인", "장르 선택", "제목 입력", "스타일 입력", "흑백 시안 생성", "시안 선택", "효과 편집", "PNG 내보내기"]
    for i, st in enumerate(steps):
        x = 0.65 + (i % 4) * 3.08
        y = 2.1 + (i // 4) * 1.45
        pill(s, x, y, 2.35, 0.72, f"{i+1}. {st}", accent=BLUE if i in {4, 6, 7} else INK)
        if i % 4 != 3:
            arrow(s, x + 2.38, y + 0.36, x + 2.75, y + 0.36, color=RGBColor(150, 164, 184))
    text_box(s, 0.95, 5.55, 11.4, 0.48, "사용자에게는 선형적이고 직관적인 제작 흐름으로 보이지만, 내부에서는 AI 생성과 이미지 처리가 단계별로 연결됩니다.", size=17, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 7)


def architecture(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "08 System Architecture", "Supabase 중심의 데이터 관리와 분리된 AI 작업 처리")
    left = [("Next.js", "사용자 화면"), ("FastAPI", "프로젝트·작업 API"), ("Worker", "AI job 처리")]
    right = [("Supabase", "Auth / DB / Private Storage"), ("OpenRouter", "LLM prompt 정리"), ("Comfy Cloud", "Diffusion generation"), ("Browser", "Effect render / export")]
    for i, (h, b) in enumerate(left):
        pill(s, 0.9, 1.75 + i * 1.1, 2.45, 0.72, h, b, accent=BLUE)
    pill(s, 5.05, 2.85, 2.7, 0.92, "Project Data", "user · project · job · asset", accent=GREEN, fill=RGBColor(235, 249, 246))
    for i, (h, b) in enumerate(right):
        pill(s, 9.65, 1.25 + i * 0.95, 2.55, 0.68, h, b, accent=NAVY)
    arrow(s, 3.4, 2.1, 5.0, 3.18)
    arrow(s, 3.4, 3.2, 5.0, 3.22)
    arrow(s, 3.4, 4.3, 5.0, 3.3)
    for yy in [1.58, 2.53, 3.48, 4.43]:
        arrow(s, 7.8, 3.25, 9.58, yy, color=BLUE)
    text_box(s, 1.0, 5.75, 11.25, 0.43, "계정 정보와 생성물은 private bucket에 저장하고, 사용자 데이터 비활용 원칙을 명시합니다.", size=16, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 8)


def current_pre_edit(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "09 Current Status - Pre Edit", "장르 선택부터 배치와 스타일 입력까지 구현")
    add_image(s, SCREENSHOTS["genre"], 0.62, 1.65, 3.0, 2.2)
    add_image(s, SCREENSHOTS["layout"], 3.9, 1.65, 3.25, 2.2)
    add_image(s, SCREENSHOTS["style"], 7.45, 1.65, 3.2, 2.2)
    text_box(s, 1.0, 4.35, 2.2, 0.24, "장르 선택", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    text_box(s, 4.55, 4.35, 2.0, 0.24, "AI 배치", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    text_box(s, 8.15, 4.35, 2.0, 0.24, "스타일 정리", size=8.5, color=MUTED, align=PP_ALIGN.CENTER)
    bullet_lines(s, 1.25, 5.15, ["장르 선택 시 내부 프롬프트 방향 설정", "제목 배치 엔진으로 기반 구도 생성", "사용자는 분위기와 요소를 자연어로 입력"], size=13.2)
    footer(s, 9)


def current_generation(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "10 Current Status - AI Generation", "3개의 흑백 타이포 시안을 만들고 선택")
    add_image(s, SCREENSHOTS["generation"], 1.05, 1.75, 5.15, 3.15)
    text_box(s, 6.85, 1.9, 5.2, 0.45, "흑백 실루엣 우선", size=20, bold=True, color=INK)
    bullet_lines(s, 6.88, 2.55, ["색·질감은 후처리 단계에서 적용", "3개 후보 중 하나를 선택하는 구조", "실패·timeout 상태를 슬롯별로 관리 예정"], size=14)
    text_box(s, 1.15, 5.65, 11.0, 0.45, "생성 모델은 형태 후보를 만들고, 사용자는 결과를 고른 뒤 편집 단계로 넘어갑니다.", size=16, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 10)


def current_effect_export(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "11 Current Status - Effect & Export", "투명화된 뼈대 위에 효과를 입히고 내보내기")
    add_image(s, SCREENSHOTS["effect"], 0.75, 1.68, 5.3, 3.2)
    add_image(s, SCREENSHOTS["export"], 6.55, 1.68, 4.7, 3.2)
    text_box(s, 1.85, 5.1, 2.8, 0.25, "효과 편집", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    text_box(s, 7.65, 5.1, 2.8, 0.25, "내보내기", size=9, color=MUTED, align=PP_ALIGN.CENTER)
    text_box(s, 0.95, 5.75, 11.3, 0.44, "보석·금속·그라데이션 효과를 웹 렌더링으로 적용하고, PNG 또는 레이어 ZIP export로 연결합니다.", size=16, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 11)


def metrics(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "12 Metric & Success", "좋은 모델보다 중요한 것은 사용자가 결과물을 받는가")
    text_box(s, 1.0, 2.0, 3.2, 0.3, "Business KPI", size=15, bold=True, color=BLUE)
    bullet_lines(s, 1.0, 2.45, ["Export 전환율", "시안 생성률", "시안 선택률", "재방문/재작업률"], size=12.8)
    text_box(s, 5.2, 2.0, 3.2, 0.3, "Model Metric", size=15, bold=True, color=PURPLE)
    bullet_lines(s, 5.2, 2.45, ["한글 가독성", "장르 적합성", "VLM side-by-side 평가", "latency / 실패율"], size=12.8, dot_color=PURPLE)
    text_box(s, 9.25, 2.0, 3.1, 0.3, "연결 방식", size=15, bold=True, color=GREEN)
    bullet_lines(s, 9.25, 2.45, ["품질 상승", "선택률 상승", "수정량 감소", "내보내기 증가"], size=12.8, dot_color=GREEN)
    text_box(s, 1.05, 5.65, 11.2, 0.43, "핵심 KPI는 생성 이후 최종 결과물을 다운로드하는 Export 전환율입니다.", size=17, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 12)


def trouble(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "13 Trouble Shooting", "가독성과 투명화 품질이 핵심 도전 과제")
    items = [("글자 가독성", "Diffusion 결과에서\n획이 흐려지거나 깨짐", RED), ("웹소설식 화려함", "장식이 많아질수록\n글자 구조가 불안정", AMBER), ("투명화 경계", "배경 제거 시\n테두리 품질 조정 필요", GREEN), ("렌더 성능", "효과 변경 시\n브라우저 비용 관리", BLUE)]
    for i, (h, b, c) in enumerate(items):
        pill(s, 0.9 + i * 3.1, 2.2, 2.55, 1.1, h, b, accent=c)
    text_box(s, 1.0, 4.65, 11.2, 0.72, "대응: 프롬프트 고도화, 흑백 실루엣 제약, threshold/곡선 기반 투명화,\n효과 렌더링 debounce와 캐시 합성 방식 적용", size=19, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 13)


def roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title(s, "14 Future Roadmap & Conclusion", "최종 발표까지는 품질과 배포 완성")
    rows = [("1", "가독성 개선", "흑백 실루엣과 한글 획 안정화"), ("2", "효과·내보내기 완성", "PNG / 레이어 ZIP 흐름 점검"), ("3", "배포 이전", "Vercel + Render + Supabase 환경 정리"), ("4", "데모 시나리오", "생성부터 export까지 끊기지 않는 시연")]
    for i, (n, h, b) in enumerate(rows):
        y = 1.95 + i * 0.95
        text_box(s, 1.05, y, 0.35, 0.28, n, size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
        text_box(s, 1.75, y, 2.8, 0.28, h, size=13.5, bold=True)
        text_box(s, 4.75, y, 6.2, 0.28, b, size=12.4, color=MUTED)
        rect(s, 1.0, y + 0.43, 10.9, 0.01, fill=RGBColor(226, 232, 240), line=RGBColor(226, 232, 240), radius=False)
    text_box(s, 1.1, 6.05, 11.0, 0.45, "작가들의 창작 활동에 실질적인 날개를 달아줄 수 있는 서비스로 발전시키겠습니다.", size=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    footer(s, 14)


def build_deck():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for fn in [
        cover, problem, product_goal, competitive, framing, pipeline, user_flow,
        architecture, current_pre_edit, current_generation, current_effect_export,
        metrics, trouble, roadmap,
    ]:
        fn(prs)
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_OUT)
    print(PPTX_OUT)


if __name__ == "__main__":
    build_deck()
